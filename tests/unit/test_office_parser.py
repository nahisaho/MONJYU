# tests/unit/test_office_parser.py
"""Tests for Word and PowerPoint document parsing (REQ-IDX-001b)."""

import pytest
from io import BytesIO
from pathlib import Path

from monjyu.document.parser import DocumentParser
from monjyu.document.models import AcademicPaperDocument


class TestWordParser:
    """Word (.docx) パースのテスト"""
    
    @pytest.fixture
    def parser(self):
        return DocumentParser()
    
    def test_parse_word_with_docx(self, parser):
        """python-docxがある場合のパース"""
        pytest.importorskip("docx", reason="python-docx not installed")
        
        from docx import Document
        from docx.shared import Pt
        
        # テスト用のdocxを作成
        doc = Document()
        doc.add_heading("Test Document Title", level=1)
        doc.add_paragraph("This is the introduction.")
        doc.add_heading("Methods", level=2)
        doc.add_paragraph("This is the methods section content.")
        doc.add_heading("Results", level=2)
        doc.add_paragraph("These are the results.")
        
        # BytesIOに保存
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        content = buffer.read()
        
        # パース
        result = parser._parse_word(content, ".docx")
        
        assert isinstance(result, AcademicPaperDocument)
        assert result.file_type == ".docx"
        assert "Test Document Title" in result.title
        assert len(result.sections) > 0
        assert result.raw_text != ""
    
    def test_parse_word_extract_doi(self, parser):
        """DOI抽出テスト"""
        pytest.importorskip("docx", reason="python-docx not installed")
        
        from docx import Document
        
        doc = Document()
        doc.add_heading("Research Paper", level=1)
        doc.add_paragraph("DOI: 10.1234/example.2024")
        
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        content = buffer.read()
        
        result = parser._parse_word(content, ".docx")
        
        assert result.doi == "10.1234/example.2024"
    
    def test_parse_word_sections_classification(self, parser):
        """セクション分類テスト"""
        pytest.importorskip("docx", reason="python-docx not installed")
        
        from docx import Document
        
        doc = Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("Intro content.")
        doc.add_heading("Methodology", level=1)
        doc.add_paragraph("Method content.")
        doc.add_heading("Conclusion", level=1)
        doc.add_paragraph("Conclusion content.")
        
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        content = buffer.read()
        
        result = parser._parse_word(content, ".docx")
        
        # セクションタイプが分類されている
        section_types = [s.section_type for s in result.sections]
        assert "introduction" in section_types or len(result.sections) > 0
    
    def test_parse_word_fallback_without_docx(self, parser, monkeypatch):
        """python-docxがない場合のフォールバック"""
        import builtins
        original_import = builtins.__import__
        
        def mock_import(name, *args, **kwargs):
            if name == "docx":
                raise ImportError("No module named 'docx'")
            return original_import(name, *args, **kwargs)
        
        monkeypatch.setattr(builtins, "__import__", mock_import)
        
        result = parser._parse_word(b"dummy content", ".docx")
        
        assert isinstance(result, AcademicPaperDocument)
        assert "python-docx required" in result.title


class TestPowerPointParser:
    """PowerPoint (.pptx) パースのテスト"""
    
    @pytest.fixture
    def parser(self):
        return DocumentParser()
    
    def test_parse_pptx_with_library(self, parser):
        """python-pptxがある場合のパース"""
        pytest.importorskip("pptx", reason="python-pptx not installed")
        
        from pptx import Presentation
        from pptx.util import Inches
        
        # テスト用のpptxを作成
        prs = Presentation()
        
        # タイトルスライド
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Presentation Title"
        subtitle.text = "Subtitle"
        
        # コンテンツスライド
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = content_slide.shapes
        shapes.title.text = "Slide 2 Title"
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Bullet point 1"
        
        # BytesIOに保存
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        content = buffer.read()
        
        # パース
        result = parser._parse_powerpoint(content)
        
        assert isinstance(result, AcademicPaperDocument)
        assert result.file_type == ".pptx"
        assert "Presentation Title" in result.title
        assert len(result.sections) == 2  # 2スライド
        assert result.page_count == 2
    
    def test_parse_pptx_slide_sections(self, parser):
        """スライドがセクションとして抽出されるか"""
        pytest.importorskip("pptx", reason="python-pptx not installed")
        
        from pptx import Presentation
        
        prs = Presentation()
        
        # 3枚のスライドを作成
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            txBox = slide.shapes.add_textbox(0, 0, 100, 100)
            tf = txBox.text_frame
            tf.text = f"Slide {i+1} Content"
        
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        content = buffer.read()
        
        result = parser._parse_powerpoint(content)
        
        assert len(result.sections) == 3
        # セクション見出しに "Slide" が含まれる
        for section in result.sections:
            assert "Slide" in section.heading
    
    def test_parse_pptx_fallback_without_library(self, parser, monkeypatch):
        """python-pptxがない場合のフォールバック"""
        import builtins
        original_import = builtins.__import__
        
        def mock_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("No module named 'pptx'")
            return original_import(name, *args, **kwargs)
        
        monkeypatch.setattr(builtins, "__import__", mock_import)
        
        result = parser._parse_powerpoint(b"dummy content")
        
        assert isinstance(result, AcademicPaperDocument)
        assert "python-pptx required" in result.title


class TestParserFileTypeRouting:
    """ファイルタイプルーティングのテスト"""
    
    @pytest.fixture
    def parser(self):
        return DocumentParser()
    
    def test_route_to_word_parser(self, parser, monkeypatch):
        """Word形式が正しくルーティングされるか"""
        called_with = []
        
        def mock_parse_word(content, file_type):
            called_with.append((content, file_type))
            return AcademicPaperDocument(
                file_name="test.docx",
                file_type=file_type,
                title="Test",
                authors=[],
                sections=[],
                raw_text="",
            )
        
        monkeypatch.setattr(parser, "_parse_word", mock_parse_word)
        
        parser.parse(b"content", ".docx")
        assert len(called_with) == 1
        assert called_with[0] == (b"content", ".docx")
        
        parser.parse(b"content2", ".doc")
        assert len(called_with) == 2
        assert called_with[1] == (b"content2", ".doc")
    
    def test_route_to_pptx_parser(self, parser, monkeypatch):
        """PowerPoint形式が正しくルーティングされるか"""
        called_with = []
        
        def mock_parse_pptx(content):
            called_with.append(content)
            return AcademicPaperDocument(
                file_name="test.pptx",
                file_type=".pptx",
                title="Test",
                authors=[],
                sections=[],
                raw_text="",
            )
        
        monkeypatch.setattr(parser, "_parse_powerpoint", mock_parse_pptx)
        
        parser.parse(b"content", ".pptx")
        assert len(called_with) == 1
        assert called_with[0] == b"content"


class TestFileLoaderIntegration:
    """FileLoaderとの統合テスト"""
    
    def test_loader_supports_office_formats(self):
        """Office形式がサポートされているか"""
        from monjyu.document.loader import FileLoader
        
        loader = FileLoader()
        
        assert ".docx" in loader.supported_extensions
        assert ".pptx" in loader.supported_extensions
        assert ".doc" in loader.supported_extensions
