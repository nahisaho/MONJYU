# Document Parser
"""
Document parsing for MONJYU.

Handles parsing of various document formats into structured AcademicPaperDocument.
Uses Strategy pattern for PDF processing (unstructured vs Azure Document Intelligence).
"""

from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING, Protocol

from monjyu.document.loader import FileLoader
from monjyu.document.models import (
    AcademicPaperDocument,
    AcademicSection,
    Author,
    ProcessedPDF,
)

if TYPE_CHECKING:
    from monjyu.document.pdf.base import PDFProcessorProtocol


class DocumentParserProtocol(Protocol):
    """ドキュメントパーサープロトコル"""
    
    def parse(self, content: bytes, file_type: str) -> AcademicPaperDocument:
        """ドキュメントをパース"""
        ...
    
    def parse_file(self, path: Path) -> AcademicPaperDocument:
        """ファイルをパース"""
        ...


class DocumentParser:
    """ドキュメントパーサー
    
    各種ドキュメント形式を AcademicPaperDocument に変換する。
    PDFの処理はStrategy patternで切り替え可能。
    
    Example:
        >>> parser = DocumentParser()
        >>> doc = parser.parse_file(Path("paper.pdf"))
    """
    
    # セクションタイプ分類用キーワード
    SECTION_KEYWORDS = {
        "introduction": ["introduction", "背景", "はじめに", "序論"],
        "methods": ["method", "approach", "methodology", "手法", "方法", "提案手法"],
        "results": ["result", "experiment", "evaluation", "結果", "実験", "評価"],
        "discussion": ["discussion", "考察", "議論", "analysis"],
        "conclusion": ["conclusion", "summary", "結論", "まとめ", "おわりに"],
        "related_work": ["related", "先行研究", "関連研究", "background"],
        "abstract": ["abstract", "概要", "要旨"],
        "references": ["references", "bibliography", "参考文献", "引用文献"],
    }
    
    def __init__(
        self,
        pdf_processor: PDFProcessorProtocol | None = None,
        file_loader: FileLoader | None = None,
    ) -> None:
        """初期化
        
        Args:
            pdf_processor: PDF処理ストラテジー
            file_loader: ファイルローダー
        """
        self._pdf_processor = pdf_processor
        self._file_loader = file_loader or FileLoader()
    
    @property
    def pdf_processor(self) -> PDFProcessorProtocol:
        """PDF処理プロセッサを取得（遅延初期化）"""
        if self._pdf_processor is None:
            from monjyu.document.pdf.unstructured_processor import UnstructuredPDFProcessor
            self._pdf_processor = UnstructuredPDFProcessor()
        return self._pdf_processor
    
    def parse(self, content: bytes, file_type: str) -> AcademicPaperDocument:
        """ドキュメントをパース
        
        Args:
            content: ファイル内容
            file_type: ファイルタイプ（拡張子）
            
        Returns:
            パースされたドキュメント
        """
        if file_type == ".pdf":
            return self._parse_pdf(content)
        elif file_type in (".txt", ".md", ".rst"):
            return self._parse_text(content, file_type)
        elif file_type in (".html", ".htm"):
            return self._parse_html(content)
        elif file_type == ".json":
            return self._parse_json(content)
        elif file_type == ".tex":
            return self._parse_latex(content)
        elif file_type in (".docx", ".doc"):
            return self._parse_word(content, file_type)
        elif file_type == ".pptx":
            return self._parse_powerpoint(content)
        else:
            # その他はテキストとして処理
            return self._parse_text(content, file_type)
    
    def parse_file(self, path: Path) -> AcademicPaperDocument:
        """ファイルをパース
        
        Args:
            path: ファイルパス
            
        Returns:
            パースされたドキュメント
        """
        file_type = self._file_loader.detect_format(path)
        content = self._file_loader.load(path)
        
        doc = self.parse(content, file_type)
        doc.file_name = path.name
        
        return doc
    
    def _parse_pdf(self, content: bytes) -> AcademicPaperDocument:
        """PDFをパース"""
        processed = self.pdf_processor.process(content)
        
        return AcademicPaperDocument(
            file_name=processed.file_name,
            file_type=".pdf",
            title=processed.title,
            authors=processed.authors,
            doi=processed.doi,
            arxiv_id=processed.arxiv_id,
            abstract=processed.abstract,
            sections=processed.sections,
            tables=processed.tables,
            figures=processed.figures,
            references=processed.references,
            keywords=processed.keywords,
            language=processed.language,
            page_count=processed.page_count,
            raw_text=processed.raw_text,
        )
    
    def _parse_text(self, content: bytes, file_type: str) -> AcademicPaperDocument:
        """テキスト系ファイルをパース"""
        text = content.decode("utf-8", errors="replace")
        
        # タイトルとセクションを抽出
        title, sections = self._extract_structure_from_text(text)
        
        return AcademicPaperDocument(
            file_name="",
            file_type=file_type,
            title=title,
            authors=[],
            sections=sections,
            raw_text=text,
        )
    
    def _parse_html(self, content: bytes) -> AcademicPaperDocument:
        """HTMLをパース"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(content, "html.parser")
            
            # タイトル抽出
            title = ""
            title_tag = soup.find("title")
            if title_tag:
                title = title_tag.get_text(strip=True)
            
            # h1タグからもタイトル候補
            h1_tag = soup.find("h1")
            if h1_tag and not title:
                title = h1_tag.get_text(strip=True)
            
            # 本文抽出
            body = soup.find("body")
            text = body.get_text(separator="\n", strip=True) if body else ""
            
            _, sections = self._extract_structure_from_text(text)
            
            return AcademicPaperDocument(
                file_name="",
                file_type=".html",
                title=title,
                authors=[],
                sections=sections,
                raw_text=text,
            )
        except ImportError:
            # BeautifulSoupがなければテキストとして処理
            return self._parse_text(content, ".html")
    
    def _parse_json(self, content: bytes) -> AcademicPaperDocument:
        """JSONをパース"""
        import json
        
        data = json.loads(content.decode("utf-8"))
        
        # 一般的な学術論文JSONスキーマを想定
        title = data.get("title", data.get("Title", ""))
        
        # 著者
        authors_data = data.get("authors", data.get("Authors", []))
        authors = []
        for a in authors_data:
            if isinstance(a, str):
                authors.append(Author(name=a))
            elif isinstance(a, dict):
                authors.append(Author(
                    name=a.get("name", a.get("Name", "")),
                    affiliation=a.get("affiliation"),
                    email=a.get("email"),
                ))
        
        abstract = data.get("abstract", data.get("Abstract", ""))
        
        return AcademicPaperDocument(
            file_name="",
            file_type=".json",
            title=title,
            authors=authors,
            abstract=abstract,
            doi=data.get("doi", data.get("DOI")),
            arxiv_id=data.get("arxiv_id", data.get("arxivId")),
            publication_year=data.get("year", data.get("Year")),
            keywords=data.get("keywords", []),
            raw_text=str(data),
        )
    
    def _parse_latex(self, content: bytes) -> AcademicPaperDocument:
        """LaTeXをパース"""
        text = content.decode("utf-8", errors="replace")
        
        # タイトル抽出
        title = ""
        title_match = re.search(r"\\title\{([^}]+)\}", text)
        if title_match:
            title = title_match.group(1)
        
        # 著者抽出
        authors = []
        author_match = re.search(r"\\author\{([^}]+)\}", text)
        if author_match:
            author_text = author_match.group(1)
            # \\and で分割
            for name in re.split(r"\\and|\\\\", author_text):
                name = re.sub(r"\\[a-z]+\{[^}]*\}", "", name)  # コマンド除去
                name = name.strip()
                if name:
                    authors.append(Author(name=name))
        
        # アブストラクト抽出
        abstract = ""
        abstract_match = re.search(
            r"\\begin\{abstract\}(.*?)\\end\{abstract\}",
            text,
            re.DOTALL
        )
        if abstract_match:
            abstract = abstract_match.group(1).strip()
        
        # セクション抽出
        sections = []
        section_pattern = r"\\section\{([^}]+)\}(.*?)(?=\\section|\\end\{document\}|$)"
        for match in re.finditer(section_pattern, text, re.DOTALL):
            heading = match.group(1)
            content_text = match.group(2).strip()
            # LaTeXコマンドを簡易除去
            content_text = re.sub(r"\\[a-z]+\{[^}]*\}", "", content_text)
            content_text = re.sub(r"\\[a-z]+", "", content_text)
            
            sections.append(AcademicSection(
                heading=heading,
                level=1,
                section_type=self.classify_section(heading),
                content=content_text,
                page_numbers=[],
            ))
        
        return AcademicPaperDocument(
            file_name="",
            file_type=".tex",
            title=title,
            authors=authors,
            abstract=abstract,
            sections=sections,
            raw_text=text,
        )
    
    def _extract_structure_from_text(
        self,
        text: str,
    ) -> tuple[str, list[AcademicSection]]:
        """テキストから構造を抽出
        
        Args:
            text: テキスト
            
        Returns:
            (タイトル, セクションリスト)
        """
        lines = text.split("\n")
        title = ""
        sections: list[AcademicSection] = []
        current_section: AcademicSection | None = None
        
        for line in lines:
            line = line.strip()
            
            if not line:
                continue
            
            # タイトル候補（最初の大きなテキスト）
            if not title and len(line) < 200 and not line.startswith("#"):
                title = line
                continue
            
            # Markdownヘッダー
            header_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if header_match:
                if current_section:
                    sections.append(current_section)
                
                level = len(header_match.group(1))
                heading = header_match.group(2)
                
                current_section = AcademicSection(
                    heading=heading,
                    level=level,
                    section_type=self.classify_section(heading),
                    content="",
                    page_numbers=[],
                )
                continue
            
            # 通常のヘッダー候補（短い行、大文字で始まる等）
            if (
                len(line) < 100
                and line[0].isupper()
                and not line.endswith(".")
                and self._is_likely_header(line)
            ):
                if current_section:
                    sections.append(current_section)
                
                current_section = AcademicSection(
                    heading=line,
                    level=1,
                    section_type=self.classify_section(line),
                    content="",
                    page_numbers=[],
                )
                continue
            
            # 本文
            if current_section:
                current_section.content += line + "\n"
        
        if current_section:
            sections.append(current_section)
        
        return title, sections
    
    def _is_likely_header(self, line: str) -> bool:
        """ヘッダーらしいか判定"""
        # 数字で始まる（番号付きセクション）
        if re.match(r"^\d+\.", line):
            return True
        
        # 全大文字
        if line.isupper():
            return True
        
        # セクションキーワードを含む
        line_lower = line.lower()
        for keywords in self.SECTION_KEYWORDS.values():
            if any(kw in line_lower for kw in keywords):
                return True
        
        return False
    
    def classify_section(self, heading: str) -> str:
        """セクションタイプを分類（IMRaD）
        
        Args:
            heading: セクション見出し
            
        Returns:
            セクションタイプ
        """
        heading_lower = heading.lower()
        
        for section_type, keywords in self.SECTION_KEYWORDS.items():
            if any(kw in heading_lower for kw in keywords):
                return section_type
        
        return "other"
    
    def extract_doi(self, text: str) -> str | None:
        """DOIを抽出
        
        Args:
            text: テキスト
            
        Returns:
            DOI文字列またはNone
        """
        doi_pattern = r"10\.\d{4,}/[^\s\]>\"']+"
        match = re.search(doi_pattern, text)
        return match.group() if match else None
    
    def extract_arxiv_id(self, text: str) -> str | None:
        """arXiv IDを抽出
        
        Args:
            text: テキスト
            
        Returns:
            arXiv ID文字列またはNone
        """
        arxiv_pattern = r"arXiv:(\d{4}\.\d{4,5}(?:v\d+)?)"
        match = re.search(arxiv_pattern, text, re.IGNORECASE)
        return match.group(1) if match else None

    def _parse_word(self, content: bytes, file_type: str) -> AcademicPaperDocument:
        """Word文書（.docx, .doc）をパース
        
        REQ-IDX-001b: Word文書の前処理
        
        Args:
            content: ファイル内容
            file_type: ファイルタイプ（.docx または .doc）
            
        Returns:
            パースされたドキュメント
        """
        try:
            from docx import Document
            
            doc = Document(BytesIO(content))
            
            # タイトル抽出（最初の段落または見出しスタイル）
            title = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    # 見出し1スタイルを優先
                    if para.style and para.style.name.startswith("Heading 1"):
                        title = para.text.strip()
                        break
                    # または最初の非空段落
                    if not title:
                        title = para.text.strip()
            
            # セクションとテキスト抽出
            sections = []
            current_heading = ""
            current_level = 1
            current_content: list[str] = []
            full_text_parts: list[str] = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                full_text_parts.append(text)
                
                # 見出しスタイルの検出
                is_heading = False
                heading_level = 1
                if para.style and para.style.name.startswith("Heading"):
                    is_heading = True
                    # "Heading 1", "Heading 2" などからレベルを抽出
                    try:
                        heading_level = int(para.style.name.split()[-1])
                    except (ValueError, IndexError):
                        heading_level = 1
                
                if is_heading:
                    # 前のセクションを保存
                    if current_heading or current_content:
                        sections.append(AcademicSection(
                            heading=current_heading,
                            level=current_level,
                            content="\n".join(current_content),
                            section_type=self.classify_section(current_heading),
                        ))
                    current_heading = text
                    current_level = heading_level
                    current_content = []
                else:
                    current_content.append(text)
            
            # 最後のセクションを保存
            if current_heading or current_content:
                sections.append(AcademicSection(
                    heading=current_heading,
                    level=current_level,
                    content="\n".join(current_content),
                    section_type=self.classify_section(current_heading),
                ))
            
            raw_text = "\n\n".join(full_text_parts)
            
            # タイトルが空の場合はraw_textの最初の行を使用
            if not title and full_text_parts:
                title = full_text_parts[0][:200]
            
            return AcademicPaperDocument(
                file_name="",
                file_type=file_type,
                title=title,
                authors=[],
                doi=self.extract_doi(raw_text),
                arxiv_id=self.extract_arxiv_id(raw_text),
                sections=sections,
                raw_text=raw_text,
            )
        
        except ImportError:
            # python-docxがインストールされていない場合
            import warnings
            warnings.warn(
                "python-docx is not installed. Install it with: pip install python-docx"
            )
            # フォールバック：バイナリとして最低限のテキスト抽出を試みる
            return AcademicPaperDocument(
                file_name="",
                file_type=file_type,
                title="[Word document - python-docx required]",
                authors=[],
                sections=[],
                raw_text="",
            )
        except Exception as e:
            # その他のパースエラー
            import warnings
            warnings.warn(f"Failed to parse Word document: {e}")
            return AcademicPaperDocument(
                file_name="",
                file_type=file_type,
                title="[Word document - parse error]",
                authors=[],
                sections=[],
                raw_text="",
            )

    def _parse_powerpoint(self, content: bytes) -> AcademicPaperDocument:
        """PowerPointプレゼンテーション（.pptx）をパース
        
        REQ-IDX-001b: PowerPoint文書の前処理
        
        Args:
            content: ファイル内容
            
        Returns:
            パースされたドキュメント
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(BytesIO(content))
            
            # タイトルスライドからタイトル抽出
            title = ""
            if prs.slides and len(prs.slides) > 0:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            title = text
                            break
            
            # 各スライドをセクションとして抽出
            sections = []
            full_text_parts: list[str] = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_title = ""
                slide_content_parts: list[str] = []
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                if not slide_title:
                                    slide_title = text
                                else:
                                    slide_content_parts.append(text)
                
                slide_content = "\n".join(slide_content_parts)
                
                if slide_title or slide_content:
                    sections.append(AcademicSection(
                        heading=f"Slide {slide_num}: {slide_title}",
                        level=1,  # スライドは全てレベル1
                        content=slide_content,
                        section_type="slide",
                    ))
                    full_text_parts.append(slide_title)
                    full_text_parts.extend(slide_content_parts)
            
            raw_text = "\n\n".join(full_text_parts)
            
            return AcademicPaperDocument(
                file_name="",
                file_type=".pptx",
                title=title or "[Untitled Presentation]",
                authors=[],
                doi=self.extract_doi(raw_text),
                arxiv_id=self.extract_arxiv_id(raw_text),
                sections=sections,
                page_count=len(prs.slides),
                raw_text=raw_text,
            )
        
        except ImportError:
            # python-pptxがインストールされていない場合
            import warnings
            warnings.warn(
                "python-pptx is not installed. Install it with: pip install python-pptx"
            )
            return AcademicPaperDocument(
                file_name="",
                file_type=".pptx",
                title="[PowerPoint - python-pptx required]",
                authors=[],
                sections=[],
                raw_text="",
            )
        except Exception as e:
            # その他のパースエラー
            import warnings
            warnings.warn(f"Failed to parse PowerPoint: {e}")
            return AcademicPaperDocument(
                file_name="",
                file_type=".pptx",
                title="[PowerPoint - parse error]",
                authors=[],
                sections=[],
                raw_text="",
            )
